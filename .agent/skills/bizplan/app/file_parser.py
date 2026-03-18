"""
file_parser.py — 업로드된 파일에서 텍스트 추출
지원 형식: .docx, .pptx, .xlsx, .txt, .md
"""

import os
from typing import Optional


def extract_text(file_path: str) -> Optional[str]:
    """파일 형식에 따라 텍스트 추출"""
    ext = os.path.splitext(file_path)[1].lower()

    extractors = {
        '.txt': _extract_txt,
        '.md': _extract_txt,
        '.docx': _extract_docx,
        '.pptx': _extract_pptx,
        '.xlsx': _extract_xlsx,
        '.pdf': _extract_pdf,
    }

    extractor = extractors.get(ext)
    if not extractor:
        return None

    try:
        return extractor(file_path)
    except Exception as e:
        return f"[파일 추출 실패: {e}]"


def _extract_pdf(path: str) -> str:
    """PDF 파일"""
    from PyPDF2 import PdfReader
    reader = PdfReader(path)
    lines = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            lines.append(text)
    return '\n'.join(lines)


def _extract_txt(path: str) -> str:
    """텍스트/마크다운 파일"""
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()


def _extract_docx(path: str) -> str:
    """Word 문서"""
    from docx import Document
    doc = Document(path)
    lines = []
    for p in doc.paragraphs:
        if p.text.strip():
            lines.append(p.text)
    # 테이블도 추출
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append(' | '.join(cells))
    return '\n'.join(lines)


def _extract_pptx(path: str) -> str:
    """PowerPoint 파일"""
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n--- 슬라이드 {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        lines.append(para.text)
    return '\n'.join(lines)


def _extract_xlsx(path: str) -> str:
    """Excel 파일"""
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    lines = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"\n--- 시트: {sheet_name} ---")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else '' for c in row]
            if any(cells):
                lines.append(' | '.join(cells))
    return '\n'.join(lines)
